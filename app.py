import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import requests

st.set_page_config(page_title="SlideToAngelo", layout="wide")

st.title("🎯 SlideToAngelo")
st.markdown("Converti testo in slide PowerPoint con stile e velocità.")

input_text = st.text_area("✍️ Inserisci il contenuto delle slide (usa `---` per separare)", height=300)

def create_slide(prs, title=None, content=None, img_url=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    if title:
        slide.shapes.title.text = title
    if content:
        slide.placeholders[1].text = content
    if img_url:
        try:
            img_data = requests.get(img_url).content
            img_stream = BytesIO(img_data)
            slide.shapes.add_picture(img_stream, Inches(5), Inches(1.5), width=Inches(4))
        except:
            st.warning(f"Impossibile caricare immagine da: {img_url}")

def parse_input(text):
    slides = text.split("---")
    parsed = []
    for slide in slides:
        lines = slide.strip().split("\n")
        title = ""
        content = ""
        img_url = None
        for line in lines:
            if line.startswith("img:"):
                img_url = line.replace("img:", "").strip()
            elif line.startswith("#"):
                title = line.replace("#", "").strip()
            else:
                content += line + "\n"
        parsed.append((title, content.strip(), img_url))
    return parsed

if st.button("🎨 Genera Slide"):
    prs = Presentation()
    slides_data = parse_input(input_text)
    for title, content, img_url in slides_data:
        create_slide(prs, title, content, img_url)

    pptx_io = BytesIO()
    prs.save(pptx_io)
    st.success("✅ Slide generate con successo!")
    st.download_button("📥 Scarica il file .pptx", data=pptx_io.getvalue(), file_name="SlideToAngelo.pptx")
